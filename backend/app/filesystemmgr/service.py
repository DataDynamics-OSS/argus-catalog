"""로컬 파일시스템 브라우저 서비스 레이어.

``os`` / ``pathlib`` 만 사용해 로컬 리눅스 파일시스템에서 디렉토리 조회·파일
업로드/다운로드·생성/삭제/이름 변경·미리보기를 구현한다. 모든 경로는
설정된 데이터 디렉토리(``settings.data_dir``) 를 기준으로 해석되며, API 가
``path="/"`` 로 호출하면 곧 ``data_dir`` 로 매핑된다.

``root_sub`` 가 주어지면 ``data_dir/<root_sub>`` 가 가상 루트가 되어 브라우저가
하위 디렉토리에 갇힌다 (예: ``model-artifacts`` / ``oci-artifacts``).

로깅 정책: 잠재적 보안 위반(``..`` / 절대 경로로 루트 탈출 시도) 은 WARNING
으로 기록한다. 변경 동작 자체의 INFO 는 호출 라우터에서 한 번만 남긴다
(이중 기록 방지).
"""

import grp
import logging
import os
import pwd
import shutil
import stat
from datetime import UTC, datetime
from pathlib import Path

from app.core.config import settings
from app.filesystemmgr.schemas import (
    CreateFolderResponse,
    DeleteResponse,
    DocumentPreviewResponse,
    FileInfo,
    FileStatResponse,
    FolderInfo,
    ListDirectoryResponse,
    RenameResponse,
    TablePreviewResponse,
)

logger = logging.getLogger(__name__)


# =========================================================================== #
# 헬퍼
# =========================================================================== #


def _get_root_dir(sub_path: str | None = None) -> Path:
    """파일 브라우저가 사용할 루트 디렉토리를 돌려준다.

    ``sub_path`` 가 주어지면 ``data_dir / sub_path`` 를 반환한다.
    예) ``sub_path="model-artifacts"`` → ``data_dir/model-artifacts``
    """
    root = settings.data_dir
    if sub_path:
        root = root / sub_path
    return root


def _format_dt(timestamp: float) -> str:
    """Unix 타임스탬프를 ISO 8601 문자열(UTC) 로 변환."""
    return datetime.fromtimestamp(timestamp, tz=UTC).isoformat()


def _permission_string(mode: int) -> str:
    """``stat`` 모드를 ``rwxrwxrwx`` 형식의 권한 문자열로 변환."""
    perms = ""
    for who in ("USR", "GRP", "OTH"):
        for what in ("R", "W", "X"):
            flag = getattr(stat, f"S_I{what}{who}")
            perms += what.lower() if mode & flag else "-"
    return perms


def _permission_octal(mode: int) -> str:
    """``stat`` 모드를 8진수 문자열(예: ``0755``) 로 변환."""
    return format(stat.S_IMODE(mode), "04o")


def _get_owner(st: os.stat_result) -> str:
    """파일의 ``uid`` 를 사용자명으로 변환. 알 수 없으면 uid 그대로 반환."""
    try:
        return pwd.getpwuid(st.st_uid).pw_name
    except KeyError:
        return str(st.st_uid)


def _get_group(st: os.stat_result) -> str:
    """파일의 ``gid`` 를 그룹명으로 변환. 알 수 없으면 gid 그대로 반환."""
    try:
        return grp.getgrgid(st.st_gid).gr_name
    except KeyError:
        return str(st.st_gid)


def _resolve_path(path: str, root_sub: str | None = None) -> Path:
    """사용자 입력 경로를 루트 디렉토리 기준으로 해석한다.

    사용자에게 ``"/"`` 는 루트로 보이며 실제로는 ``data_dir`` 또는
    ``data_dir/<root_sub>`` 에 매핑된다.
    예) ``data_dir=/var/lib/argus-catalog-server``, ``root_sub="model-artifacts"`` 일 때:
      - ``path="/"``      → ``/var/lib/argus-catalog-server/model-artifacts``
      - ``path="/iris"``  → ``/var/lib/argus-catalog-server/model-artifacts/iris``

    ``..`` 등을 통해 루트를 탈출하려는 경로는 ``ValueError`` 로 차단한다(WARNING 로깅).
    """
    root = _get_root_dir(root_sub).resolve()
    root.mkdir(parents=True, exist_ok=True)

    relative = path.lstrip("/")
    if relative:
        resolved = (root / relative).resolve()
    else:
        resolved = root

    if not (resolved == root or str(resolved).startswith(str(root) + "/")):
        logger.warning("경로 탈출 시도: input=%s, resolved=%s, root=%s", path, resolved, root)
        raise ValueError("Access denied: path escapes data directory")

    return resolved


def _to_user_path(resolved: Path, root_sub: str | None = None) -> str:
    """내부 해석 경로를 사용자에게 표시되는 가상 경로로 환원한다."""
    root = _get_root_dir(root_sub).resolve()
    try:
        rel = resolved.relative_to(root)
        return "/" + str(rel) if str(rel) != "." else "/"
    except ValueError:
        return str(resolved)


# =========================================================================== #
# 1. 디렉토리 목록
# =========================================================================== #


async def list_directory(path: str, root_sub: str | None = None) -> ListDirectoryResponse:
    """주어진 경로 아래의 파일/디렉토리 목록 + 메타데이터를 돌려준다."""
    resolved = _resolve_path(path, root_sub)

    if not resolved.exists():
        raise FileNotFoundError(f"Directory not found: {resolved}")
    if not resolved.is_dir():
        raise NotADirectoryError(f"Not a directory: {resolved}")

    folders: list[FolderInfo] = []
    files: list[FileInfo] = []

    try:
        entries = sorted(resolved.iterdir(), key=lambda e: e.name)
    except PermissionError:
        logger.warning("디렉토리 접근 권한 없음: %s", resolved)
        raise PermissionError(f"Permission denied: {resolved}")

    for entry in entries:
        try:
            st = entry.stat(follow_symlinks=False)
        except (PermissionError, OSError):
            continue

        owner = _get_owner(st)
        group = _get_group(st)
        perms = _permission_string(st.st_mode)

        if stat.S_ISDIR(st.st_mode):
            folders.append(FolderInfo(
                key=_to_user_path(entry, root_sub) + "/",
                name=entry.name,
                owner=owner,
                group=group,
                permissions=perms,
            ))
        else:
            files.append(FileInfo(
                key=_to_user_path(entry, root_sub),
                name=entry.name,
                size=st.st_size,
                last_modified=_format_dt(st.st_mtime),
                owner=owner,
                group=group,
                permissions=perms,
            ))

    return ListDirectoryResponse(
        folders=folders,
        files=files,
        current_path=_to_user_path(resolved, root_sub),
    )


# =========================================================================== #
# 2. 디렉토리 생성
# =========================================================================== #


async def create_folder(path: str, root_sub: str | None = None) -> CreateFolderResponse:
    """새 디렉토리를 생성한다. 상위 디렉토리가 없으면 함께 만든다."""
    resolved = _resolve_path(path, root_sub)
    resolved.mkdir(parents=True, exist_ok=True)
    logger.info("폴더 생성: %s", resolved)
    return CreateFolderResponse(path=_to_user_path(resolved, root_sub))


# =========================================================================== #
# 3. 파일/디렉토리 삭제
# =========================================================================== #


async def delete_paths(paths: list[str], root_sub: str | None = None) -> DeleteResponse:
    """파일/디렉토리를 일괄 삭제. 디렉토리는 ``shutil.rmtree`` 로 재귀 삭제."""
    deleted: list[str] = []
    errors: list[dict] = []

    for p in paths:
        try:
            resolved = _resolve_path(p, root_sub)
            if not resolved.exists():
                errors.append({"path": p, "error": "Not found"})
                continue

            if resolved.is_dir():
                shutil.rmtree(resolved)
            else:
                resolved.unlink()

            deleted.append(p)
            logger.info("삭제됨: %s", resolved)
        except Exception as e:
            errors.append({"path": p, "error": str(e)})
            logger.error("삭제 오류: %s - %s", p, e)

    return DeleteResponse(deleted=deleted, errors=errors)


# =========================================================================== #
# 4. 이름 변경 / 이동
# =========================================================================== #


async def rename(source_path: str, destination_path: str, root_sub: str | None = None) -> RenameResponse:
    """파일/디렉토리 이름 변경 또는 이동(같은 ``os.rename`` 호출로 처리)."""
    src = _resolve_path(source_path, root_sub)
    dst = _resolve_path(destination_path, root_sub)

    if not src.exists():
        raise FileNotFoundError(f"Source not found: {src}")
    if dst.exists():
        raise FileExistsError(f"Destination already exists: {dst}")

    src.rename(dst)
    logger.info("이름 변경: %s -> %s", src, dst)
    return RenameResponse(source=_to_user_path(src, root_sub), destination=_to_user_path(dst, root_sub))


# =========================================================================== #
# 5. 파일 메타데이터 (stat)
# =========================================================================== #


async def file_stat(path: str, root_sub: str | None = None) -> FileStatResponse:
    """파일/디렉토리의 상세 메타데이터(크기·권한·시간 등) 를 ``stat`` 으로 조회."""
    resolved = _resolve_path(path, root_sub)

    if not resolved.exists():
        raise FileNotFoundError(f"Not found: {resolved}")

    st = resolved.stat(follow_symlinks=False)
    symlink_target = str(resolved.readlink()) if resolved.is_symlink() else None

    return FileStatResponse(
        path=_to_user_path(resolved, root_sub),
        name=resolved.name or "/",
        is_directory=resolved.is_dir(),
        size=st.st_size,
        last_modified=_format_dt(st.st_mtime),
        last_accessed=_format_dt(st.st_atime),
        created=_format_dt(st.st_ctime),
        owner=_get_owner(st),
        group=_get_group(st),
        permissions=_permission_string(st.st_mode),
        permissions_octal=_permission_octal(st.st_mode),
        inode=st.st_ino,
        hard_links=st.st_nlink,
        symlink_target=symlink_target,
    )


# =========================================================================== #
# 6. 파일 다운로드 (바이트 읽기)
# =========================================================================== #


async def read_file(path: str, root_sub: str | None = None) -> tuple[bytes, str]:
    """파일 내용을 읽어 ``(bytes, filename)`` 으로 돌려준다.

    라우터가 스트리밍 응답을 구성할 수 있도록 원본 바이트를 반환한다.
    """
    resolved = _resolve_path(path, root_sub)

    if not resolved.exists():
        raise FileNotFoundError(f"File not found: {resolved}")
    if resolved.is_dir():
        raise IsADirectoryError(f"Cannot download a directory: {resolved}")

    data = resolved.read_bytes()
    return data, resolved.name


# =========================================================================== #
# 7. 파일 업로드
# =========================================================================== #


async def save_uploaded_file(
    destination_dir: str,
    filename: str,
    content: bytes,
    root_sub: str | None = None,
) -> str:
    """업로드된 파일을 지정 디렉토리에 저장. 대상 디렉토리가 없으면 자동 생성."""
    dir_path = _resolve_path(destination_dir, root_sub)
    if not dir_path.is_dir():
        raise NotADirectoryError(f"Not a directory: {dir_path}")

    file_path = dir_path / filename
    file_path.write_bytes(content)
    logger.info("업로드: %s (%d bytes)", file_path, len(content))
    return _to_user_path(file_path, root_sub)


# =========================================================================== #
# 8. 파일 미리보기
# =========================================================================== #

MAX_PREVIEW_ROWS = 1000


def _serialize_value(val):
    """셀 값을 JSON 직렬화 가능한 타입으로 변환(datetime, Decimal, bytes 등)."""
    if val is None:
        return None
    if isinstance(val, (int, float, str, bool)):
        return val
    if isinstance(val, datetime):
        return val.isoformat()
    if isinstance(val, bytes):
        return val.hex()
    return str(val)


async def preview_parquet(path: str, max_rows: int = MAX_PREVIEW_ROWS, root_sub: str | None = None) -> TablePreviewResponse:
    """PyArrow 로 Parquet 파일을 읽어 테이블 형태로 미리보기."""
    import io

    import pyarrow.parquet as pq

    resolved = _resolve_path(path, root_sub)
    data = resolved.read_bytes()
    pf = pq.ParquetFile(io.BytesIO(data))
    total_rows = pf.metadata.num_rows
    columns = pf.schema_arrow.names

    table = pf.read_row_groups(list(range(pf.metadata.num_row_groups)))
    if table.num_rows > max_rows:
        table = table.slice(0, max_rows)

    rows = []
    for batch in table.to_batches():
        cols = [batch.column(i).to_pylist() for i in range(batch.num_columns)]
        for row_idx in range(batch.num_rows):
            rows.append([_serialize_value(cols[ci][row_idx]) for ci in range(len(cols))])

    logger.info("Parquet 미리보기: %s rows=%d/%d", resolved, len(rows), total_rows)
    return TablePreviewResponse(
        format="parquet",
        columns=columns,
        rows=rows,
        total_rows=total_rows,
    )


async def preview_xlsx(
    path: str,
    sheet: str | None = None,
    max_rows: int = MAX_PREVIEW_ROWS,
    root_sub: str | None = None,
) -> TablePreviewResponse:
    """openpyxl 로 XLSX/XLS 파일을 읽어 시트 단위 테이블로 미리보기."""
    import io

    import openpyxl

    resolved = _resolve_path(path, root_sub)
    data = resolved.read_bytes()
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)

    sheet_names = wb.sheetnames
    active_sheet = sheet if sheet and sheet in sheet_names else sheet_names[0]
    ws = wb[active_sheet]

    rows: list[list] = []
    columns: list[str] = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i == 0:
            columns = [str(c) if c is not None else f"col_{j}" for j, c in enumerate(row)]
            continue
        if i > max_rows:
            break
        rows.append([_serialize_value(c) for c in row])

    total_rows = ws.max_row - 1 if ws.max_row else 0
    wb.close()

    ext = resolved.suffix.lstrip(".").lower() or "xlsx"
    logger.info("Xlsx 미리보기: %s sheet=%s rows=%d/%d", resolved, active_sheet, len(rows), total_rows)
    return TablePreviewResponse(
        format=ext,
        columns=columns,
        rows=rows,
        total_rows=total_rows,
        sheet_names=sheet_names,
        active_sheet=active_sheet,
    )


async def preview_docx(path: str, root_sub: str | None = None) -> DocumentPreviewResponse:
    """mammoth 로 DOCX 파일을 HTML 로 변환해 문서 형태로 미리보기."""
    import io

    import mammoth

    resolved = _resolve_path(path, root_sub)
    data = resolved.read_bytes()
    result = mammoth.convert_to_html(io.BytesIO(data))
    if result.messages:
        logger.warning("Docx 미리보기 경고: %s", result.messages)

    logger.info("Docx 미리보기: %s html_len=%d", resolved, len(result.value))
    return DocumentPreviewResponse(format="docx", html=result.value)


async def preview_pptx(path: str, root_sub: str | None = None) -> DocumentPreviewResponse:
    """python-pptx 로 PPTX 의 슬라이드 텍스트·노트를 추출해 미리보기."""
    import io

    from pptx import Presentation

    resolved = _resolve_path(path, root_sub)
    data = resolved.read_bytes()
    prs = Presentation(io.BytesIO(data))

    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(row_texts))

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "slide_number": i,
            "texts": texts,
            "notes": notes,
        })

    html_parts = []
    for s in slides:
        html_parts.append(f'<div class="slide"><h3>Slide {s["slide_number"]}</h3>')
        for t in s["texts"]:
            html_parts.append(f"<p>{t}</p>")
        if s["notes"]:
            html_parts.append(f'<blockquote class="notes">{s["notes"]}</blockquote>')
        html_parts.append("</div><hr/>")

    logger.info("Pptx 미리보기: %s slides=%d", resolved, len(slides))
    return DocumentPreviewResponse(
        format="pptx",
        html="\n".join(html_parts),
        slides=slides,
    )
